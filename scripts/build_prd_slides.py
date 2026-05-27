"""Build docs/PRD_EXPLAINER.pptx from docs/DEMO_SCRIPT.md.

Rerun after editing the script:
    python scripts/build_prd_slides.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

REPO = Path(__file__).resolve().parent.parent
OUT = REPO / "docs" / "PRD_EXPLAINER.pptx"

NAVY = RGBColor(0x0B, 0x2A, 0x4A)
INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x55, 0x60, 0x70)
ACCENT = RGBColor(0xA8, 0x36, 0x1A)
PAPER = RGBColor(0xFA, 0xF7, 0xF2)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
RULE = RGBColor(0xE2, 0xDD, 0xD3)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_rect(slide, x, y, w, h, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    set_fill(shape, color)
    shape.shadow.inherit = False
    return shape


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_page_chrome(slide, page_num, total, section_label):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, PAPER)
    add_rect(slide, Inches(0.5), Inches(0.55), Inches(0.18), Inches(0.18), ACCENT)
    add_text(slide, Inches(0.78), Inches(0.45), Inches(8), Inches(0.4),
             section_label, size=12, bold=True, color=MUTED)
    add_text(slide, Inches(11.5), Inches(0.45), Inches(1.4), Inches(0.4),
             f"{page_num} / {total}", size=11, color=MUTED, align=PP_ALIGN.RIGHT)
    add_rect(slide, Inches(0.5), Inches(6.85), Inches(12.3), Pt(1.2), RULE)
    add_text(slide, Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.35),
             "PRD Explainer — Commercial Property Inspections",
             size=10, color=MUTED)


def add_title(slide, title, subtitle=None):
    add_text(slide, Inches(0.5), Inches(0.95), Inches(12.3), Inches(0.9),
             title, size=40, bold=True, color=NAVY)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(1.85), Inches(12.3), Inches(0.5),
                 subtitle, size=18, color=MUTED)


def add_card(slide, x, y, w, h, *, heading=None, body=None, kicker=None,
             heading_color=NAVY):
    add_rect(slide, x, y, w, h, CARD)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.08))
    set_fill(bar, ACCENT)
    pad = Inches(0.3)
    cursor_y = y + Inches(0.25)
    if kicker:
        add_text(slide, x + pad, cursor_y, w - 2 * pad, Inches(0.35),
                 kicker, size=11, bold=True, color=ACCENT)
        cursor_y += Inches(0.4)
    if heading:
        add_text(slide, x + pad, cursor_y, w - 2 * pad, Inches(0.6),
                 heading, size=22, bold=True, color=heading_color)
        cursor_y += Inches(0.7)
    if body:
        add_text(slide, x + pad, cursor_y, w - 2 * pad, y + h - cursor_y - Inches(0.2),
                 body, size=14, color=INK)


def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    add_rect(slide, Inches(0.5), Inches(0.6), Inches(0.18), Inches(0.18), ACCENT)
    add_text(slide, Inches(0.78), Inches(0.5), Inches(8), Inches(0.4),
             "PRD EXPLAINER  ·  ≤ 5 MIN", size=12, bold=True,
             color=RGBColor(0xE2, 0xDD, 0xD3))
    add_text(slide, Inches(0.5), Inches(2.2), Inches(12.3), Inches(1.6),
             "Commercial Property\nInspections",
             size=60, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    add_text(slide, Inches(0.5), Inches(4.7), Inches(12.3), Inches(0.6),
             "Replacing a legacy field-inspection workflow with a browser app.",
             size=20, color=RGBColor(0xE2, 0xDD, 0xD3))
    add_text(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.5),
             "Problem  ·  Users  ·  Must-haves",
             size=14, color=ACCENT, bold=True)
    add_text(slide, Inches(11), Inches(6.5), Inches(2), Inches(0.5),
             "reppenp", size=14, color=RGBColor(0xE2, 0xDD, 0xD3),
             align=PP_ALIGN.RIGHT)


def slide_agenda(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_chrome(slide, page, total, "AGENDA")
    add_title(slide, "What this video covers",
              "Three sections, ~5 minutes total. No code, no app demo.")
    items = [
        ("01", "The problem", "Legacy tools, no tracking, no cycle-time data."),
        ("02", "The users",  "Jeff (admin), John (inspector), Kelly (underwriter)."),
        ("03", "Must-haves", "Five user stories that define v1."),
    ]
    card_w = Inches(4.0)
    gap = Inches(0.15)
    total_w = card_w * len(items) + gap * (len(items) - 1)
    start_x = (SLIDE_W - total_w) / 2
    for i, (num, head, body) in enumerate(items):
        x = start_x + (card_w + gap) * i
        add_card(slide, x, Inches(2.9), card_w, Inches(3.3),
                 kicker=num, heading=head, body=body)


def slide_problem_intro(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_chrome(slide, page, total, "01 · THE PROBLEM")
    add_title(slide, "The team can't see its own work.",
              "Commercial property inspections, run on tooling nobody can maintain.")
    add_text(slide, Inches(0.5), Inches(2.6), Inches(12.3), Inches(0.5),
             "Three concrete pains:", size=16, bold=True, color=MUTED)
    pains = [
        ("No browser submission.",
         "Inspectors drive back to the office to type up every report. Half a day of windshield time, every inspection."),
        ("No workflow tracking.",
         "Once an order leaves admin's hands, it's a black box. Phone calls and Post-its until someone says it's done."),
        ("No cycle-time data.",
         "Nobody can tell you whether an inspection takes three days or three weeks. Can't improve what you can't measure."),
    ]
    card_w = Inches(4.0)
    gap = Inches(0.15)
    total_w = card_w * len(pains) + gap * (len(pains) - 1)
    start_x = (SLIDE_W - total_w) / 2
    for i, (head, body) in enumerate(pains):
        x = start_x + (card_w + gap) * i
        add_card(slide, x, Inches(3.2), card_w, Inches(3.3),
                 kicker=f"PAIN {i + 1}", heading=head, body=body)


def slide_problem_metric(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_chrome(slide, page, total, "01 · THE PROBLEM")
    add_title(slide, "Why this project exists",
              "Cycle time is the metric the legacy system can't produce.")
    box_w = Inches(8.5)
    box_h = Inches(3.4)
    x = (SLIDE_W - box_w) / 2
    y = Inches(2.8)
    add_rect(slide, x, y, box_w, box_h, CARD)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, box_w, Inches(0.08))
    set_fill(bar, ACCENT)
    add_text(slide, x + Inches(0.4), y + Inches(0.3), box_w - Inches(0.8), Inches(0.5),
             "Order created  →  Decision recorded", size=14, bold=True, color=MUTED)
    add_text(slide, x + Inches(0.4), y + Inches(0.9), box_w - Inches(0.8), Inches(1.4),
             "How many days?", size=44, bold=True, color=NAVY)
    add_text(slide, x + Inches(0.4), y + Inches(2.3), box_w - Inches(0.8), Inches(0.9),
             "Nobody on the team can answer this today. Surfacing it on a shared\n"
             "dashboard is the headline goal of v1.",
             size=16, color=INK)


def slide_users_overview(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_chrome(slide, page, total, "02 · THE USERS")
    add_title(slide, "Three internal users, one workflow",
              "No policyholder access in v1. Auth is a hard prerequisite for v2.")
    users = [
        ("JEFF", "Admin / Operations",
         "Creates inspection orders from policy and claims reports. Wants to know where every order stands at a glance.",
         "Desktop · Office or home"),
        ("JOHN", "Inspector",
         "Walks the building, takes notes and photos. Loses cell signal in the field. ~5 inspections/week — up to 5/day in catastrophe events.",
         "Laptop in the field · Phone is a stretch goal"),
        ("KELLY", "Underwriter",
         "Reviews submitted inspections, records premium direction and policy action. Works from a queue.",
         "Desktop · Office"),
    ]
    card_w = Inches(4.0)
    gap = Inches(0.15)
    total_w = card_w * len(users) + gap * (len(users) - 1)
    start_x = (SLIDE_W - total_w) / 2
    for i, (name, role, body, device) in enumerate(users):
        x = start_x + (card_w + gap) * i
        h = Inches(3.6)
        add_rect(slide, x, Inches(2.7), card_w, h, CARD)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(2.7), card_w, Inches(0.08))
        set_fill(bar, ACCENT)
        pad = Inches(0.3)
        add_text(slide, x + pad, Inches(2.95), card_w - 2 * pad, Inches(0.4),
                 name, size=12, bold=True, color=ACCENT)
        add_text(slide, x + pad, Inches(3.35), card_w - 2 * pad, Inches(0.5),
                 role, size=20, bold=True, color=NAVY)
        add_text(slide, x + pad, Inches(3.9), card_w - 2 * pad, Inches(1.9),
                 body, size=14, color=INK)
        add_text(slide, x + pad, Inches(5.8), card_w - 2 * pad, Inches(0.5),
                 device, size=12, bold=True, color=MUTED)


def slide_must_haves(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_chrome(slide, page, total, "03 · MUST-HAVES")
    add_title(slide, "Five must-haves for v1",
              "Direct from PRD §4. Everything else is v2 or out of scope.")
    items = [
        ("01", "Jeff creates inspection orders.",
         "Insured details, property address & use, scheduling contact. Without this, nothing else happens."),
        ("02", "John completes a 4-section inspection form in the browser.",
         "Saves mid-inspection — connectivity is unreliable. Required fields enforced at submit, DB stays permissive."),
        ("03", "Kelly reviews submissions and records a decision.",
         "Notes + premium direction (increase / decrease / no change) + policy action (approve / cancel / renew)."),
        ("04", "Cycle time is visible to everyone.",
         "Order created → decision recorded, in days, on the shared dashboard. The success metric for v1."),
        ("05", "All three users see every order's status.",
         "One shared dashboard. Ordered → In Progress → Submitted → Reviewed. No more 'where is that report?' calls."),
    ]
    top = Inches(2.5)
    row_h = Inches(0.82)
    gap = Inches(0.1)
    x = Inches(0.5)
    w = SLIDE_W - Inches(1.0)
    for i, (num, head, body) in enumerate(items):
        y = top + (row_h + gap) * i
        add_rect(slide, x, y, w, row_h, CARD)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.1), row_h)
        set_fill(bar, ACCENT)
        add_text(slide, x + Inches(0.3), y + Inches(0.1), Inches(0.6), Inches(0.6),
                 num, size=20, bold=True, color=ACCENT, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, x + Inches(1.0), y + Inches(0.08), Inches(5.5), Inches(0.4),
                 head, size=15, bold=True, color=NAVY)
        add_text(slide, x + Inches(1.0), y + Inches(0.4), w - Inches(1.2), Inches(0.45),
                 body, size=12, color=INK)


def slide_not_in_v1(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_chrome(slide, page, total, "03 · MUST-HAVES")
    add_title(slide, "Deliberately not in v1",
              "Cutting scope this hard is the only reason a three-week build is plausible.")
    items = [
        ("Policyholder portal", "Internal users only. External access is a v2+ conversation."),
        ("Text / email alerts", "Dashboard polling is enough for a controlled v1 pilot."),
        ("CSV import from policy or claims", "Jeff types orders manually for now. CSV is v2."),
        ("Authentication", "Open access in v1. Hard prerequisite before real policyholder data."),
        ("Mobile-optimized UI", "Laptop in the field is the commitment. Phone is a stretch goal."),
        ("Analytics beyond cycle time", "One metric, done well, beats five half-built dashboards."),
    ]
    cols = 3
    rows = 2
    card_w = Inches(4.0)
    card_h = Inches(1.8)
    gap = Inches(0.15)
    total_w = card_w * cols + gap * (cols - 1)
    start_x = (SLIDE_W - total_w) / 2
    start_y = Inches(2.8)
    for idx, (head, body) in enumerate(items):
        r = idx // cols
        c = idx % cols
        x = start_x + (card_w + gap) * c
        y = start_y + (card_h + gap) * r
        add_rect(slide, x, y, card_w, card_h, CARD)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.08), card_h)
        set_fill(bar, MUTED)
        add_text(slide, x + Inches(0.25), y + Inches(0.2), card_w - Inches(0.4), Inches(0.5),
                 head, size=16, bold=True, color=NAVY)
        add_text(slide, x + Inches(0.25), y + Inches(0.7), card_w - Inches(0.4), card_h - Inches(0.8),
                 body, size=12, color=INK)


def slide_close(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    add_text(slide, Inches(0.5), Inches(0.5), Inches(12.3), Inches(0.4),
             "CLOSE", size=12, bold=True, color=ACCENT)
    add_text(slide, Inches(0.5), Inches(2.4), Inches(12.3), Inches(1.4),
             "One problem.\nThree users.\nFive must-haves.",
             size=52, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    add_text(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.6),
             "Demo video shows it actually working.",
             size=22, color=RGBColor(0xE2, 0xDD, 0xD3))
    add_text(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.5),
             "ai-coding-preppen.pjreppen.workers.dev", size=14, color=ACCENT, bold=True)
    add_text(slide, Inches(11), Inches(6.5), Inches(2), Inches(0.5),
             f"{page} / {total}", size=11, color=RGBColor(0xE2, 0xDD, 0xD3),
             align=PP_ALIGN.RIGHT)


# ─── Appendix: architecture + ERD ──────────────────────────────────────────


def slide_appendix_divider(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    add_text(slide, Inches(0.5), Inches(0.5), Inches(12.3), Inches(0.4),
             "APPENDIX", size=12, bold=True, color=ACCENT)
    add_text(slide, Inches(0.5), Inches(2.4), Inches(12.3), Inches(1.4),
             "How it's built", size=60, bold=True,
             color=RGBColor(0xFF, 0xFF, 0xFF))
    add_text(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.6),
             "Cloudflare Workers + D1 + R2.",
             size=22, color=RGBColor(0xE2, 0xDD, 0xD3))
    add_text(slide, Inches(0.5), Inches(5.1), Inches(12.3), Inches(0.6),
             "Architecture and data model in two slides.",
             size=18, color=RGBColor(0xE2, 0xDD, 0xD3))
    add_text(slide, Inches(11), Inches(6.5), Inches(2), Inches(0.5),
             f"{page} / {total}", size=11, color=RGBColor(0xE2, 0xDD, 0xD3),
             align=PP_ALIGN.RIGHT)


def _box(slide, x, y, w, h, *, fill=CARD, border=None):
    rect = add_rect(slide, x, y, w, h, fill)
    if border:
        rect.line.color.rgb = border
        rect.line.width = Pt(0.75)
    return rect


def _label_box(slide, x, y, w, h, *, title, subtitle=None,
               title_size=13, subtitle_size=10):
    _box(slide, x, y, w, h)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.06))
    set_fill(bar, ACCENT)
    pad = Inches(0.12)
    add_text(slide, x + pad, y + Inches(0.14), w - 2 * pad, Inches(0.32),
             title, size=title_size, bold=True, color=NAVY)
    if subtitle:
        add_text(slide, x + pad, y + Inches(0.46), w - 2 * pad,
                 h - Inches(0.5), subtitle, size=subtitle_size, color=INK)


def _down_arrow(slide, center_x, y, h=Inches(0.18), w=Inches(0.28)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW, center_x - w / 2, y, w, h,
    )
    set_fill(shape, MUTED)


def slide_architecture(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_chrome(slide, page, total, "APPENDIX · ARCHITECTURE")
    add_title(slide, "Request flow at a glance",
              "Three users → one SPA → one Worker → split storage (D1 + R2).")

    center_x = SLIDE_W / 2

    # Row 1 — three user pills
    user_w = Inches(2.0)
    user_h = Inches(0.5)
    gap = Inches(0.3)
    row_w = user_w * 3 + gap * 2
    start_x = (SLIDE_W - row_w) / 2
    users_y = Inches(2.05)
    for i, name in enumerate(["Jeff · Admin", "John · Inspector", "Kelly · Underwriter"]):
        x = start_x + (user_w + gap) * i
        pill = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, users_y, user_w, user_h,
        )
        set_fill(pill, NAVY)
        add_text(slide, x, users_y, user_w, user_h, name, size=12, bold=True,
                 color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)

    _down_arrow(slide, center_x, Inches(2.65))

    # Row 2 — SPA (one wide box)
    band_x = Inches(1.0)
    band_w = SLIDE_W - Inches(2.0)
    _label_box(
        slide, band_x, Inches(2.92), band_w, Inches(0.55),
        title="React SPA",
        subtitle="routes: /  ·  /orders/new  ·  /orders/:id  ·  /review  ·  /review/:id     (served by ASSETS binding · SPA fallback)",
        title_size=13, subtitle_size=9,
    )

    _down_arrow(slide, center_x, Inches(3.55))

    # Row 3 — Worker
    _label_box(
        slide, band_x, Inches(3.78), band_w, Inches(0.55),
        title="Hono Worker · src/index.ts",
        subtitle="GET /health  ·  app.route('/api/orders', …) → orders · forms · photos · decisions",
        title_size=13, subtitle_size=9,
    )

    _down_arrow(slide, center_x, Inches(4.4))

    # Row 4 — four route modules
    routes = [
        ("orders.ts", "POST · /\nGET · / (cycle_time_days)"),
        ("forms.ts", "GET/PUT · /:id/form\nPOST · /:id/submit"),
        ("photos.ts", "POST · /:id/photos\nGET · /:id/photos[/:photoId]"),
        ("decisions.ts", "POST · /:id/decision\nGET · /:id/decision"),
    ]
    route_w = Inches(2.95)
    route_h = Inches(0.85)
    route_gap = Inches(0.12)
    route_row_w = route_w * 4 + route_gap * 3
    route_start = (SLIDE_W - route_row_w) / 2
    routes_y = Inches(4.65)
    for i, (name, paths) in enumerate(routes):
        x = route_start + (route_w + route_gap) * i
        _label_box(slide, x, routes_y, route_w, route_h, title=name,
                   subtitle=paths, title_size=12, subtitle_size=9)

    _down_arrow(slide, center_x, Inches(5.6))

    # Row 5 — storage (D1 with 4 tables · R2 with image objects)
    storage_y = Inches(5.85)
    storage_h = Inches(0.95)
    d1_w = Inches(7.6)
    r2_w = Inches(4.0)
    d1_x = Inches(0.5)
    r2_x = SLIDE_W - r2_w - Inches(0.5)

    _label_box(slide, d1_x, storage_y, d1_w, storage_h,
               title="D1 · binding DB · database 'inspections'",
               title_size=12)
    tbl_y = storage_y + Inches(0.4)
    tbl_h = Inches(0.45)
    tbl_w = (d1_w - Inches(0.4) - Inches(0.21)) / 4
    for i, t in enumerate(["inspections", "form_responses", "photos", "decisions"]):
        tx = d1_x + Inches(0.15) + (tbl_w + Inches(0.07)) * i
        _box(slide, tx, tbl_y, tbl_w, tbl_h, fill=PAPER, border=RULE)
        add_text(slide, tx, tbl_y, tbl_w, tbl_h, t, size=10, bold=True,
                 color=NAVY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    _label_box(slide, r2_x, storage_y, r2_w, storage_h,
               title="R2 · binding PHOTOS · bucket 'inspection-photos'",
               title_size=12)
    obj_y = storage_y + Inches(0.4)
    obj_w = r2_w - Inches(0.3)
    _box(slide, r2_x + Inches(0.15), obj_y, obj_w, tbl_h, fill=PAPER, border=RULE)
    add_text(slide, r2_x + Inches(0.15), obj_y, obj_w, tbl_h,
             "image objects (r2_key)", size=10, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def _entity_box(slide, x, y, w, h, *, name, attrs):
    """ERD entity: navy header + attribute lines on a white card."""
    _box(slide, x, y, w, h, fill=CARD, border=RULE)
    header_h = Inches(0.42)
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, header_h)
    set_fill(header, NAVY)
    add_text(slide, x, y, w, header_h, name, size=14, bold=True,
             color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    body_x = x + Inches(0.18)
    body_y = y + header_h + Inches(0.08)
    body_w = w - Inches(0.36)
    body_h = h - header_h - Inches(0.16)
    add_text(slide, body_x, body_y, body_w, body_h,
             "\n".join(f"• {a}" for a in attrs),
             size=10, color=INK)


def slide_erd(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_page_chrome(slide, page, total, "APPENDIX · ERD")
    add_title(slide, "Data model (D1)",
              "inspections is the spine. Every other table joins by inspection_id and cascades on delete.")

    # Left: INSPECTIONS (tall, full vertical height of the diagram area).
    ins_x = Inches(0.5)
    ins_y = Inches(2.1)
    ins_w = Inches(5.2)
    ins_h = Inches(4.6)
    _entity_box(slide, ins_x, ins_y, ins_w, ins_h, name="INSPECTIONS",
                attrs=[
                    "id (PK)",
                    "status — Ordered / In Progress / Submitted / Reviewed",
                    "source — Policy / Claims (nullable)",
                    "insured_name (NOT NULL)",
                    "property_address (NOT NULL)",
                    "property_use",
                    "contact_name · contact_phone",
                    "assigned_inspector",
                    "created_at · updated_at",
                    "submitted_at — set on POST /:id/submit",
                    "decided_at — set on POST /:id/decision",
                    "    (drives cycle_time_days on GET /api/orders)",
                ])

    # Right: three siblings stacked.
    right_x = Inches(7.6)
    right_w = Inches(5.2)
    sib_h = Inches(1.45)

    fr_y = Inches(2.1)
    _entity_box(slide, right_x, fr_y, right_w, sib_h, name="FORM_RESPONSES",
                attrs=[
                    "id (PK)  ·  inspection_id (FK · UNIQUE)",
                    "Section 1 (10 cols)  ·  Section 2 (11 cols)",
                    "Section 3 (6 cols)  ·  Section 4 (7 cols)",
                    "updated_at — all answer columns NULLABLE",
                ])

    dec_y = Inches(3.7)
    _entity_box(slide, right_x, dec_y, right_w, sib_h, name="DECISIONS",
                attrs=[
                    "id (PK)  ·  inspection_id (FK · UNIQUE)",
                    "premium_direction — increase / decrease / no change",
                    "policy_action — approve / cancel / renew",
                    "notes  ·  decided_by  ·  created_at",
                ])

    ph_y = Inches(5.3)
    _entity_box(slide, right_x, ph_y, right_w, sib_h, name="PHOTOS",
                attrs=[
                    "id (PK)  ·  inspection_id (FK)",
                    "section — CHECK 1..4",
                    "r2_key (UNIQUE) — pointer to R2 object",
                    "filename · content_type · size_bytes · uploaded_at",
                ])

    # Relationship lines + cardinality labels.
    line_x_start = ins_x + ins_w
    line_x_end = right_x
    for sib_top, cardinality in [
        (fr_y, "1  ──  0..1"),
        (dec_y, "1  ──  0..1"),
        (ph_y, "1  ──  0..N"),
    ]:
        mid_y = sib_top + sib_h / 2
        line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, line_x_start, mid_y, line_x_end, mid_y,
        )
        line.line.color.rgb = MUTED
        line.line.width = Pt(1.5)
        # White-filled label so the line doesn't strike through the text.
        lbl_w = Inches(1.4)
        lbl_h = Inches(0.28)
        lbl_x = (line_x_start + line_x_end) / 2 - lbl_w / 2
        lbl_y = mid_y - lbl_h / 2
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, lbl_x, lbl_y, lbl_w, lbl_h)
        set_fill(bg, PAPER)
        add_text(slide, lbl_x, lbl_y, lbl_w, lbl_h, cardinality,
                 size=10, bold=True, color=MUTED, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        slide_title,
        slide_agenda,
        slide_problem_intro,
        slide_problem_metric,
        slide_users_overview,
        slide_must_haves,
        slide_not_in_v1,
        slide_close,
        slide_appendix_divider,
        slide_architecture,
        slide_erd,
    ]
    total = len(builders)
    for i, builder in enumerate(builders, start=1):
        if builder is slide_title:
            builder(prs)
        else:
            builder(prs, i, total)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Wrote {OUT.relative_to(REPO)} ({total} slides)")


if __name__ == "__main__":
    build()
